import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from tkinter import messagebox, filedialog

def convert_pdf_to_ppt():
    try:
        # select pdf
        input_path = filedialog.askopenfilename(
            title="Select PDF file",
            filetypes=[("PDF Files", "*.pdf")]
        )
        if not input_path:
            return
        
        # select output folder
        output_dir = filedialog.askdirectory(title="Select Output folder")
        if not output_dir:
            return
        
        pdf_name = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{pdf_name}.pptx")

        # convert pdf to images
        images = convert_from_path(input_path)

        # create PowerPoint presentation
        prs = Presentation()

        # set slide dimensions to 16:9 aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for img in images:
            # create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank slide layout
            img_path = os.path.join(output_dir, "temp_page.png")
            img.save(img_path, "PNG")

            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

            os.remove(img_path)  # remove temporary image file

        prs.save(output_path)

        messagebox.showinfo("Success", f"PDF converted to PowerPoint successfully!\nSaved at: {output_path}")

    except Exception as e:
        messagebox.showerror("Error", f"Failed to convert PDF to PowerPoint:\n\n{str(e)}")